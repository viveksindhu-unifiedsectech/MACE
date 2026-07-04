#!/usr/bin/env python3
"""Generate the Secure Files patent addendum (.docx) and slide deck (.pptx).
Run with a Python that has python-docx and python-pptx installed."""
import os
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PColor

HERE = os.path.dirname(os.path.abspath(__file__))
NAVY = RGBColor(0x0A, 0x1F, 0x44)
PNAVY = PColor(0x0A, 0x1F, 0x44)
PWHITE = PColor(0xFF, 0xFF, 0xFF)
PACCENT = PColor(0x1E, 0x90, 0xFF)

# ─────────────────────────── PATENT ADDENDUM ───────────────────────────
CLAIMS = [
    ("31. File-type-agnostic envelope encryption with cryptographic tenant binding",
     "A method for securing data at rest wherein each file, regardless of type, is "
     "encrypted under a unique per-file data encryption key (DEK), the DEK being wrapped "
     "by a hardware-security-module-backed master key, and wherein the wrap operation is "
     "bound to an authenticated encryption context comprising a tenant identifier such that "
     "a wrapped DEK issued for a first tenant cannot be unwrapped under a second tenant's "
     "context, thereby enforcing tenant isolation cryptographically rather than by access "
     "policy alone."),
    ("32. Chunked authenticated container resisting reorder, truncation, and splice",
     "The method of claim 31 wherein the ciphertext is organized as an ordered sequence of "
     "independently authenticated chunks, each chunk authenticated with associated data "
     "derived from a digest of the file header and the chunk's ordinal position and an "
     "end-of-stream marker, such that removal, reordering, duplication, or substitution of "
     "any chunk is detected on decryption."),
    ("33. Pre-encryption redaction with tamper-evident report",
     "A method wherein sensitive content, comprising personally identifiable information and "
     "leaked secrets, is detected and removed from a file before encryption and storage, and "
     "wherein a report enumerating the categories and counts of removed items is produced and "
     "recorded in an immutable audit log without recording the removed values themselves."),
    ("34. Pre-action AI safeguard",
     "A method wherein, prior to completing a data action comprising upload, share, or "
     "download, an automated safeguard evaluates the action and its payload to produce a risk "
     "score and a verdict of allow, warn, or block, the evaluation combining deterministic "
     "detectors with an optional language-model assessment operating on a redaction-safe "
     "summary from which raw sensitive values have been removed, and wherein a critical "
     "verdict prevents the action from completing."),
    ("35. Layered access control with non-overridable tenant isolation",
     "An access-control method combining role-based and attribute-based authorization with a "
     "data-classification clearance gate, wherein tenant isolation is evaluated as an "
     "outermost, categorical condition that no role, including an administrative role, may "
     "override, and wherein a named-subject grant may raise a subject's effective clearance "
     "for a single object without altering the subject's clearance for any other object."),
    ("36. Privacy-preserving cross-matter conflict-of-interest detection",
     "A method for detecting conflicts of interest across segregated matters wherein entities "
     "extracted from documents are represented only as keyed one-way tokens computed with a "
     "tenant-scoped secret, and wherein the appearance of a common token across two or more "
     "matters separated by an ethical wall is reported as a potential conflict, the report "
     "referencing the token and matter identifiers but not any raw entity value, such that the "
     "detection index cannot itself disclose the underlying data."),
    ("37. Privilege-leak detection via cross-privilege token co-occurrence",
     "The method of claim 36 wherein a keyed entity token appearing in both a document marked "
     "privileged and a document not so marked is reported as a potential privilege waiver."),
    ("38. Tenant-scoped token non-comparability",
     "The method of claim 36 wherein the keyed token for a given entity value differs across "
     "tenants, rendering tokens non-comparable between tenants and reinforcing tenant isolation "
     "at the correlation layer."),
    ("39. Defense-in-depth double encryption",
     "The method of claim 31 wherein the application-layer envelope ciphertext is additionally "
     "encrypted by an independent server-side encryption facility of the object store keyed by a "
     "distinct managed key, such that compromise of either layer alone does not disclose "
     "plaintext."),
    ("40. Unified secure-ingest pipeline with immutable audit",
     "A system comprising the safeguard of claim 34, the redaction of claim 33, and the envelope "
     "encryption of claim 31 arranged as a single ingest pipeline that guards, then redacts, then "
     "encrypts, then stores each file, and that records each action in an append-only audit log "
     "sufficient for SOC 2, FedRAMP, and GDPR evidentiary requirements."),
]


def build_patent():
    doc = Document()
    t = doc.add_heading("MACE — Patent Application Addendum", level=0)
    p = doc.add_paragraph("Secure File Encryption, Access Control, Redaction, "
                          "AI Safeguard & Cross-Matter Conflict Detection")
    p.runs[0].bold = True
    doc.add_paragraph("UnifiedSec Technologies · Patent IN/2026/UNISEC/MACE-001 + PCT")
    doc.add_paragraph("Addendum claims 31–40 — supplements the 30 claims of the base filing. "
                      "DRAFT for counsel review; not yet filed.")

    doc.add_heading("Field of the addendum", level=1)
    doc.add_paragraph(
        "This addendum extends the MACE platform from threat correlation to securing data at "
        "rest. It claims a unified pipeline that encrypts files of any type with per-file, "
        "HSM-wrapped keys bound to tenant identity; controls access by role, attribute, and "
        "classification under non-overridable tenant isolation; redacts sensitive content before "
        "encryption; warns of risk before an action completes; and detects conflicts of interest "
        "and privilege leaks across segregated matters using a privacy-preserving keyed-token "
        "index that stores no raw data.")

    doc.add_heading("Claims", level=1)
    for title, body in CLAIMS:
        h = doc.add_paragraph()
        r = h.add_run(title); r.bold = True; r.font.color.rgb = NAVY
        doc.add_paragraph(body)

    doc.add_heading("Reduction to practice", level=1)
    doc.add_paragraph(
        "The claimed methods are implemented in mace_platform/backend/app/secure/ and exercised "
        "by an automated test suite (tests/test_secure_files.py) and an end-to-end demonstration "
        "(scripts/secure_files_demo.py) that runs offline, confirming: envelope encryption with "
        "tenant-bound context; chunk-level tamper detection; redaction before encryption; a "
        "safeguard that blocks a leaked private key; classification and tenant-isolation denials; "
        "and cross-matter conflict and privilege-leak detection over a keyed-token index proven "
        "to hold no raw values.")
    out = os.path.join(HERE, "MACE_Patent_Addendum_SecureFiles.docx")
    doc.save(out)
    return out


# ─────────────────────────── SLIDE DECK ───────────────────────────
SLIDES = [
    ("MACE Secure Files",
     "The one-stop cybersecurity layer for data at rest\nUnifiedSec Technologies · Patent-pending"),
    ("The problem",
     "• Organizations hold files full of secrets & PII with weak, inconsistent protection\n"
     "• Access is by folder permission, not by identity + classification\n"
     "• No one warns you BEFORE a risky upload or over-broad share\n"
     "• Nothing catches the same party showing up across walled matters"),
    ("The vision: one-stop, universal",
     "Encrypt ANY file → push to AWS securely → only the right people open it →\n"
     "redact what shouldn't leak → AI warns before harm → detect conflicts across matters.\n"
     "Works for law firms, banks, healthcare, government — any regulated data."),
    ("How it works — the secure-ingest pipeline",
     "AI GUARD  →  REDACT  →  ENVELOPE-ENCRYPT  →  AWS S3 (SSE-KMS)\n"
     "Every step logged to an immutable audit trail (SOC2 / FedRAMP / GDPR)."),
    ("Encryption you can defend",
     "• Per-file 256-bit key, wrapped by an AWS KMS (HSM) master key\n"
     "• Key wrap bound to tenant identity → cryptographic tenant isolation\n"
     "• Chunked authentication → tamper, truncation & reorder are detected\n"
     "• Defense in depth: our envelope + AWS SSE-KMS (two independent layers)"),
    ("Access by identity, role & classification",
     "• RBAC + ABAC + data classification (public→restricted)\n"
     "• Tenant isolation is categorical — no admin overrides it\n"
     "• Named-user grants unlock a single file without widening anyone's access\n"
     "• Every allow/deny is explainable and audited"),
    ("Redaction + AI safeguard",
     "• Python engine strips SSNs, cards, keys, tokens BEFORE encryption\n"
     "• AI guard scores each action and BLOCKS a leaked private key up front\n"
     "• Flags over-broad shares, executables, and prompt-injection attempts\n"
     "• Optional Claude second opinion — on redaction-safe summaries only"),
    ("The differentiator: cross-matter conflict detection",
     "• Correlates entities across documents & matters using keyed hashes\n"
     "• Catches conflicts of interest across ethical walls\n"
     "• Detects privilege leaks (privileged entity in a public filing)\n"
     "• Privacy-preserving: the index stores NO raw client data\n"
     "No commodity DLP/encryption tool does this — it is the patent thesis."),
    ("AWS cloud security",
     "• KMS-wrapped keys, S3 block-public + versioning, TLS-only, KMS-only writes\n"
     "• Least-privilege IAM, CloudTrail audit, multi-region data residency\n"
     "• Terraform + Helm + Docker Compose — runs local for demos, EKS for scale"),
    ("Proven, tested, demoable today",
     "• 24 automated tests pass offline (crypto, access, redaction, guard, conflict)\n"
     "• One command runs the full end-to-end demo on a laptop — no AWS needed\n"
     "• Docker Compose brings up API + Elasticsearch + Kibana for a live server"),
    ("Honest roadmap to production",
     "Done: encryption, access control, redaction, AI guard, conflict engine, tests, IaC.\n"
     "Before real client data: external pen test, third-party code audit, SOC 2 Type I.\n"
     "These are named in the threat model — no surprises for a buyer's security team."),
    ("Patent-pending · UnifiedSec Technologies",
     "10 new claims (31–40) extend MACE to data-at-rest security.\n"
     "Contact: sindhuvick8@gmail.com"),
]


def build_deck():
    prs = Presentation()
    prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for i, (title, body) in enumerate(SLIDES):
        s = prs.slides.add_slide(blank)
        bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PNAVY
        # accent bar
        bar = s.shapes.add_shape(1, PInches(0), PInches(0), PInches(0.25), PInches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = PACCENT; bar.line.fill.background()
        # title
        tb = s.shapes.add_textbox(PInches(0.8), PInches(0.6), PInches(11.8), PInches(1.4))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = PPt(40 if i == 0 else 32); r.font.bold = True; r.font.color.rgb = PWHITE
        # body
        bb = s.shapes.add_textbox(PInches(0.9), PInches(2.1), PInches(11.6), PInches(4.8))
        bf = bb.text_frame; bf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            para = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            rr = para.add_run(); rr.text = line
            rr.font.size = PPt(22 if i == 0 else 18)
            rr.font.color.rgb = PACCENT if (i == 0 and j == 0) else PWHITE
        # footer
        fb = s.shapes.add_textbox(PInches(0.9), PInches(6.95), PInches(11.6), PInches(0.4))
        fr = fb.text_frame.paragraphs[0].add_run()
        fr.text = "UnifiedSec MACE · Secure Files · Patent-pending IN/2026/UNISEC/MACE-001"
        fr.font.size = PPt(10); fr.font.color.rgb = PColor(0x88, 0x99, 0xBB)
    out = os.path.join(HERE, "MACE_SecureFiles_Deck_v1.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    print("patent:", build_patent())
    print("deck:  ", build_deck())
